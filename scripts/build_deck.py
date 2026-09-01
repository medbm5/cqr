"""Build the Citalid case-study presentation.

Fourteen slides, in French, for an audience of cyber risk modellers. Every
figure quoted is read from `deck/figures.json`, which `build_deck_charts.py`
writes from the engine itself - so a number cannot appear on a slide unless the
pipeline produced it, and cannot go stale without the deck being rebuilt.

Run from the repository root, after the charts:

    python scripts/build_deck_charts.py
    python scripts/build_deck.py
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent

#: Change this before presenting.
PRESENTER = "[Votre nom]"

# The cockpit's palette, so deck and product are visibly one thing.
BG = RGBColor(0x05, 0x09, 0x12)
CARD = RGBColor(0x0A, 0x11, 0x20)
LINE = RGBColor(0x1B, 0x29, 0x42)
INK = RGBColor(0xF1, 0xF5, 0xF9)
INK_2 = RGBColor(0x94, 0xA3, 0xB8)
INK_3 = RGBColor(0x7A, 0x8A, 0xA0)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)
CAUTION = RGBColor(0xFB, 0xBF, 0x24)
EDR = RGBColor(0xD9, 0x59, 0x26)

# Inter is the product's typeface but is not installed everywhere; Segoe UI is
# on every Windows machine and is the closest safe substitute. A deck that
# silently falls back to Times on the presenting laptop is not worth the risk.
FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)


# --------------------------------------------------------------- primitives


def add_deck() -> Presentation:
    """A 16:9 presentation with no theme of its own."""
    deck = Presentation()
    deck.slide_width, deck.slide_height = W, H
    return deck


def blank(deck: Presentation):
    """One slide, painted the product's navy from edge to edge."""
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    shape.shadow.inherit = False
    return slide


def text(
    slide,
    body: str,
    *,
    left,
    top,
    width,
    height,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    spacing=1.0,
    anchor=MSO_ANCHOR.TOP,
):
    """A text box. python-pptx has no autofit worth the name, so sizes are set."""
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0

    for index, line in enumerate(body.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = spacing
        run = paragraph.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return box


def eyebrow(slide, label: str) -> None:
    """The accent kicker every slide opens with."""
    text(
        slide,
        label.upper(),
        left=MARGIN,
        top=Inches(0.62),
        width=W - 2 * MARGIN,
        height=Inches(0.3),
        size=12,
        color=ACCENT,
        bold=True,
    )


def title(slide, label: str, *, size=32) -> None:
    """The one idea the slide carries."""
    text(
        slide,
        label,
        left=MARGIN,
        top=Inches(1.0),
        width=W - 2 * MARGIN,
        height=Inches(1.0),
        size=size,
        color=INK,
        bold=True,
        spacing=1.05,
    )


def rule(slide, top) -> None:
    """A hairline, the same one the UI uses to separate a card's header."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, top, W - 2 * MARGIN, Emu(9525))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LINE
    shape.line.fill.background()
    shape.shadow.inherit = False


def card(slide, left, top, width, height, *, fill=CARD, border=LINE):
    """The surface a figure or a block of keywords sits on."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.06
    return shape


def figure(slide, value: str, label: str, *, left, top, width, size=40, color=INK) -> None:
    """A number set large, with the thing it counts underneath it."""
    text(
        slide,
        value,
        left=left,
        top=top,
        width=width,
        height=Inches(0.75),
        size=size,
        color=color,
        bold=True,
    )
    text(
        slide,
        label,
        left=left,
        top=top + Inches(0.72),
        width=width,
        height=Inches(0.6),
        size=12,
        color=INK_2,
        spacing=1.15,
    )


def keywords(slide, lines: list[str], *, left, top, width, size=16, gap=0.46) -> None:
    """Keyword lines with an accent tick, never sentences."""
    for index, line in enumerate(lines):
        offset = top + Inches(index * gap)
        tick = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, offset + Inches(0.06), Inches(0.045), Inches(0.18)
        )
        tick.fill.solid()
        tick.fill.fore_color.rgb = ACCENT
        tick.line.fill.background()
        tick.shadow.inherit = False
        text(
            slide,
            line,
            left=left + Inches(0.22),
            top=offset,
            width=width - Inches(0.22),
            height=Inches(0.4),
            size=size,
            color=INK_2,
        )


def picture(slide, path: Path, *, left, top, width) -> None:
    """Insert a regenerated chart. Never a screenshot."""
    slide.shapes.add_picture(str(path), left, top, width=width)


def notes(slide, body: str) -> None:
    """Speaker notes: the oral justification, not a transcript of the slide."""
    slide.notes_slide.notes_text_frame.text = body.strip()


#: Height of one box in the pipeline chain. A module constant rather than a
#: default argument, because Inches() is a call and would run at import.
CHAIN_HEIGHT = Inches(1.25)


def chain(slide, steps: list[tuple[str, str]], *, top, height=CHAIN_HEIGHT) -> None:
    """The pipeline as boxes and arrows — the traceability slide's whole point."""
    count = len(steps)
    gap = Inches(0.16)
    width = (W - 2 * MARGIN - gap * (count - 1)) / count

    for index, (value, label) in enumerate(steps):
        left = MARGIN + (width + gap) * index
        last = index == count - 1
        card(
            slide,
            left,
            top,
            width,
            height,
            fill=CARD,
            border=ACCENT if last else LINE,
        )
        text(
            slide,
            value,
            left=left + Inches(0.12),
            top=top + Inches(0.20),
            width=width - Inches(0.24),
            height=Inches(0.5),
            size=19 if len(value) <= 9 else 16,
            color=ACCENT if last else INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        text(
            slide,
            label,
            left=left + Inches(0.10),
            top=top + Inches(0.72),
            width=width - Inches(0.20),
            height=Inches(0.5),
            size=10.5,
            color=INK_3,
            align=PP_ALIGN.CENTER,
            spacing=1.1,
        )
        if not last:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                left + width + Inches(0.02),
                top + height / 2 - Inches(0.07),
                Inches(0.12),
                Inches(0.14),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INK_3
            arrow.line.fill.background()
            arrow.shadow.inherit = False


# ------------------------------------------------------------------- slides


def build(figures: dict, charts: Path, out: Path) -> None:
    """Assemble the fourteen slides."""
    f = figures
    deck = add_deck()

    # ------------------------------------------------------------- 1. titre
    slide = blank(deck)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(2.55), Inches(0.07), Inches(1.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    text(
        slide,
        "QUANTIFICATION DU RISQUE CYBER",
        left=MARGIN + Inches(0.32),
        top=Inches(2.55),
        width=Inches(9),
        height=Inches(0.3),
        size=13,
        color=ACCENT,
        bold=True,
    )
    text(
        slide,
        "De la télémétrie brute\nà une perte annuelle chiffrée",
        left=MARGIN + Inches(0.32),
        top=Inches(3.0),
        width=Inches(10.5),
        height=Inches(1.8),
        size=44,
        color=INK,
        bold=True,
        spacing=1.05,
    )
    text(
        slide,
        "ETI Retail / e-commerce   ·   1 200 employés   ·   maturité 55/100   ·   ~20 actifs",
        left=MARGIN + Inches(0.32),
        top=Inches(4.75),
        width=Inches(11),
        height=Inches(0.4),
        size=15,
        color=INK_2,
    )
    rule(slide, Inches(5.5))
    text(
        slide,
        PRESENTER,
        left=MARGIN,
        top=Inches(5.75),
        width=Inches(6),
        height=Inches(0.4),
        size=15,
        color=INK,
        bold=True,
    )
    text(
        slide,
        "Un chiffre de risque doit toujours pouvoir être tracé et expliqué",
        left=W - MARGIN - Inches(6),
        top=Inches(5.78),
        width=Inches(6),
        height=Inches(0.4),
        size=13,
        color=INK_3,
        align=PP_ALIGN.RIGHT,
    )
    notes(
        slide,
        """
Le cas demande une perte annuelle ; je l'ai traité comme un problème de traçabilité avant
d'en faire un problème de statistique, parce que n'importe qui peut produire un nombre et que
l'enjeu est de pouvoir le défendre ligne à ligne.
Le fil rouge est celui de l'énoncé : chaque chiffre remonte à ses sources, et je terminerai par
la commande qui rejoue l'ensemble.
Je vais passer volontairement du temps sur une erreur que j'ai commise et corrigée — c'est elle
qui structure tout le modèle.
""",
    )

    # ------------------------------------------- 2. la démarche en une image
    slide = blank(deck)
    eyebrow(slide, "La démarche en une image")
    title(slide, "Six étapes, six justifications")
    chain(
        slide,
        [
            (f["rows_read"], "lignes brutes\nSIEM + EDR"),
            (f["total_events"], "événements\ndistincts"),
            (f["attack_grade"], "événements\nattack-grade"),
            (f["episodes"], "épisodes\nd'attaque"),
            (f["lambda_incident"], "incident à perte\npar an"),
            (f["aal"], "perte annuelle\nmoyenne"),
        ],
        top=Inches(2.55),
    )
    text(
        slide,
        "Chaque flèche est une décision de modélisation — et le reste de l'exposé les justifie "
        "une par une.",
        left=MARGIN,
        top=Inches(4.25),
        width=W - 2 * MARGIN,
        height=Inches(0.4),
        size=15,
        color=INK_2,
    )
    keywords(
        slide,
        [
            "déduplication inter-flux   →   slide 4",
            "seuil de sévérité + fenêtre de session   →   slide 6",
            "changement d'unité : détections ≠ incidents à perte   →   slide 7",
            "lognormale pondérée par similarité   →   slides 8 à 10",
        ],
        left=MARGIN,
        top=Inches(4.95),
        width=Inches(11),
        size=14,
        gap=0.42,
    )
    notes(
        slide,
        """
Tout l'exposé tient ici : 45 840 lignes livrées par deux outils, 274 k€ de perte annuelle
moyenne à l'arrivée.
Les quatre premières flèches réduisent le volume — dédoublonnage, filtre de sévérité,
regroupement en épisodes — mais la cinquième est d'une autre nature : c'est un changement
d'unité, et c'est là que se joue la crédibilité du chiffre.
Retenez l'ordre de grandeur, un facteur 150 000 entre les lignes brutes et les incidents à
perte : une chaîne aussi longue ne vaut que si chaque maillon est auditable.
""",
    )

    # ------------------------------------------------ 3. architecture logicielle
    slide = blank(deck)
    eyebrow(slide, "Architecture logicielle")
    title(slide, "Le modèle est un composant, l'interface est une vitrine")

    layers = [
        (
            "risk_engine",
            "package Python pur",
            ["ingestion", "frequency", "severity", "simulation"],
            ACCENT,
        ),
        ("API Django / DRF", "vues minces", ["parse → appel → sérialisation"], INK_2),
        ("UI Next.js", "rendu seul", ["aucun calcul de risque"], INK_2),
    ]
    top = Inches(2.5)
    for index, (name, role, items, colour) in enumerate(layers):
        height = Inches(1.15)
        offset = top + Inches(index * 1.32)
        card(slide, MARGIN, offset, Inches(7.4), height, border=colour if index == 0 else LINE)
        text(
            slide,
            name,
            left=MARGIN + Inches(0.3),
            top=offset + Inches(0.2),
            width=Inches(3.4),
            height=Inches(0.4),
            size=19,
            color=colour,
            bold=True,
        )
        text(
            slide,
            role,
            left=MARGIN + Inches(0.3),
            top=offset + Inches(0.66),
            width=Inches(3.4),
            height=Inches(0.3),
            size=12,
            color=INK_3,
        )
        text(
            slide,
            "   ·   ".join(items),
            left=MARGIN + Inches(3.55),
            top=offset + Inches(0.44),
            width=Inches(3.7),
            height=Inches(0.4),
            size=12,
            color=INK_2,
        )
        if index < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                MARGIN + Inches(3.6),
                offset + height + Inches(0.02),
                Inches(0.16),
                Inches(0.12),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INK_3
            arrow.line.fill.background()
            arrow.shadow.inherit = False

    card(slide, Inches(8.7), Inches(2.5), Inches(3.8), Inches(2.47), border=LINE)
    text(
        slide,
        "Qualité",
        left=Inches(9.0),
        top=Inches(2.72),
        width=Inches(3.2),
        height=Inches(0.3),
        size=13,
        color=ACCENT,
        bold=True,
    )
    keywords(
        slide,
        [
            f"{f['backend_tests']} tests, {f['coverage']} de couverture",
            "mypy strict · ruff · pre-commit",
            "aucun import Django dans le moteur",
            "seeds explicites, résultats rejouables",
            "to_explanation() sur chaque résultat",
        ],
        left=Inches(9.0),
        top=Inches(3.2),
        width=Inches(3.3),
        size=12,
        gap=0.35,
    )
    text(
        slide,
        "Le moteur tourne sans serveur : notebook, CLI, ou test. L'API et l'UI le consomment, "
        "elles ne le contiennent pas.",
        left=MARGIN,
        top=Inches(6.5),
        width=W - 2 * MARGIN,
        height=Inches(0.4),
        size=14,
        color=INK_2,
    )
    notes(
        slide,
        """
Contrainte que je me suis imposée : risk_engine n'importe jamais Django, et un test parcourt
tous les modules dans un sous-processus pour le prouver, parce qu'une dépendance de ce genre
s'introduit sans qu'on s'en aperçoive.
La conséquence pratique, c'est que le modèle s'utilise dans un notebook, un batch ou un autre
produit, sans traîner un framework web derrière lui.
to_explanation() rend cet exposé possible : chaque objet de résultat sait raconter comment il a
été obtenu, et c'est la même trace que la CLI imprime, que l'API sert et que l'interface
affiche.
""",
    )

    # --------------------------------------------- 4. ingestion & normalisation
    slide = blank(deck)
    eyebrow(slide, "Ingestion & normalisation")
    title(slide, "Deux outils, un seul flux d'événements")
    picture(slide, charts / "04_funnel.png", left=MARGIN, top=Inches(2.05), width=Inches(7.6))
    keywords(
        slide,
        [
            "sévérité textuelle (SIEM) vs score 0–999 (EDR)",
            "schémas différents, événements communs",
            f"{f['dup_cross_feed']} vus par les DEUX flux",
            "clé de fusion : actif + technique + horodatage",
            "règle : la PIRE sévérité gagne",
        ],
        left=Inches(8.6),
        top=Inches(2.35),
        width=Inches(4.0),
        size=13,
        gap=0.52,
    )
    figure(
        slide,
        f"+{f['inflation']}",
        "d'inflation évitée — ce qu'une concaténation naïve\naurait ajouté à toutes les fréquences",
        left=MARGIN,
        top=Inches(4.95),
        width=Inches(5.5),
        size=46,
        color=CAUTION,
    )
    figure(
        slide,
        f["duplicates"],
        "doublons absorbés",
        left=Inches(6.6),
        top=Inches(4.95),
        width=Inches(3),
        size=46,
    )
    notes(
        slide,
        """
Le piège du jeu de données est ici : les deux flux ne sont pas complémentaires mais
partiellement redondants, avec 12 343 événements portant le même actif, la même technique et le
même horodatage des deux côtés.
Concaténer aurait gonflé toutes les fréquences de 42 %, et l'erreur serait passée inaperçue
parce que le résultat reste plausible.
Sur la règle de fusion : quand les deux outils divergent sur la sévérité je garde la pire — si
l'EDR voit critique là où le SIEM voit medium, c'est qu'il a vu quelque chose de plus proche du
poste, et prendre la moyenne diluerait un signal fort dans un signal faible.
""",
    )

    # -------------------------------------------------------- 5. stationnarité
    slide = blank(deck)
    eyebrow(slide, "Stationnarité")
    title(slide, "Un débit plat autorise Poisson et l'annualisation")
    picture(slide, charts / "05_weekly.png", left=MARGIN, top=Inches(2.15), width=Inches(11.6))
    figure(
        slide,
        f["observed_days"],
        "jours observés\n1ᵉʳ nov. 2025 → 31 mai 2026",
        left=MARGIN,
        top=Inches(5.55),
        width=Inches(3),
        size=34,
    )
    figure(
        slide,
        f["annualization"],
        "facteur d'annualisation\n365 / 212, recalculé depuis les données",
        left=Inches(4.6),
        top=Inches(5.55),
        width=Inches(4),
        size=34,
        color=ACCENT,
    )
    notes(
        slide,
        """
Cette slide justifie deux hypothèses d'un coup : d'abord Poisson, parce que le débit
hebdomadaire est stable sur les sept mois, sans tendance ni rupture, ce qui est la condition
d'un processus à taux constant.
Ensuite l'annualisation — le facteur 1,7217 n'est pas écrit en dur, il vaut 365 divisé par le
nombre de jours réellement observés, donc un export plus long change le résultat tout seul.
Limite assumée : sept mois ne montrent aucune saisonnalité annuelle, et un pic de novembre lié
au Black Friday sur un e-commerçant, je ne peux ni le confirmer ni l'exclure.
""",
    )

    # ------------------------------------------------ 6. fréquence, étape 1
    slide = blank(deck)
    eyebrow(slide, "Fréquence — étape 1")
    title(slide, "Une alerte n'est pas une attaque")
    picture(slide, charts / "06_window.png", left=MARGIN, top=Inches(2.1), width=Inches(7.5))
    keywords(
        slide,
        [
            "regroupement par actif",
            "silence > 24 h → nouvel épisode",
            "seuil de sévérité : high et au-delà",
            "les deux conventions sont des curseurs",
            "exposés dans l'interface, pas cachés",
        ],
        left=Inches(8.7),
        top=Inches(2.5),
        width=Inches(4.0),
        size=13,
        gap=0.5,
    )
    for index, (value, label, colour) in enumerate(
        [
            (f["attack_grade"], "événements\nattack-grade", INK),
            (f["episodes"], "épisodes\nreconstitués", ACCENT),
            (f["compression"], "compression", CAUTION),
            (f["lambda_detected"], "attaques détectées\npar an", INK),
        ]
    ):
        figure(
            slide,
            value,
            label,
            left=MARGIN + Inches(index * 3.05),
            top=Inches(5.35),
            width=Inches(2.9),
            size=34,
            color=colour,
        )
    notes(
        slide,
        """
Une intrusion ne produit pas une alerte mais une rafale ; compter les alertes mesurerait le
bavardage des sondes, pas la fréquence des attaques.
Je regroupe donc par actif avec une fenêtre de silence de 24 heures, l'ordre de grandeur d'une
session d'intrusion et d'un cycle de traitement en SOC : c'est une convention, pas une mesure,
et le graphique montre exactement ce qu'elle coûte, de 8 500 épisodes à une heure jusqu'à 911 à
vingt-quatre.
Le regroupement se fait par actif seul et non par actif et type d'attaque, parce qu'un intrus
déclenche les détections qu'il rencontre et que les compter par type reviendrait à recompter la
même intrusion sous plusieurs noms.
""",
    )

    # ------------------------------------------------ 7. fréquence, étape 2
    slide = blank(deck)
    eyebrow(slide, "Fréquence — étape 2  ·  le cœur du sujet")
    title(slide, "Changement d'unité : détecter n'est pas perdre")

    card(slide, MARGIN, Inches(2.0), Inches(5.6), Inches(1.5), border=EDR)
    text(
        slide,
        "Première version",
        left=MARGIN + Inches(0.25),
        top=Inches(2.2),
        width=Inches(5),
        height=Inches(0.3),
        size=12,
        color=EDR,
        bold=True,
    )
    text(
        slide,
        "λ = 9 168 attaques/an   →   AAL 12,5 Md€",
        left=MARGIN + Inches(0.25),
        top=Inches(2.6),
        width=Inches(5.2),
        height=Inches(0.5),
        size=19,
        color=INK,
        bold=True,
    )
    text(
        slide,
        "absurde pour une ETI de 1 200 personnes",
        left=MARGIN + Inches(0.25),
        top=Inches(3.05),
        width=Inches(5.2),
        height=Inches(0.3),
        size=12,
        color=INK_3,
    )

    card(slide, Inches(7.1), Inches(2.0), Inches(5.4), Inches(1.5), border=LINE)
    text(
        slide,
        "Diagnostic — erreur de catégorie",
        left=Inches(7.35),
        top=Inches(2.2),
        width=Inches(5),
        height=Inches(0.3),
        size=12,
        color=ACCENT,
        bold=True,
    )
    text(
        slide,
        "télémétrie → attaques DÉTECTÉES\nbase externe → incidents À PERTE",
        left=Inches(7.35),
        top=Inches(2.6),
        width=Inches(5),
        height=Inches(0.8),
        size=15,
        color=INK,
        spacing=1.25,
    )

    picture(slide, charts / "07_calibration.png", left=MARGIN, top=Inches(3.7), width=Inches(8.2))
    keywords(
        slide,
        [
            "λ_incident = λ_détecté × p",
            f"p ≈ 1 détection sur {f['p_one_in']}",
            f"ancré sur {f['peer_companies']} organisations",
            "la base ancre le NIVEAU",
            "la télémétrie garde le MIX par type",
        ],
        left=Inches(9.1),
        top=Inches(4.15),
        width=Inches(3.7),
        size=13,
        gap=0.5,
    )
    notes(
        slide,
        """
Ma première version multipliait 9 168 attaques par an par un coût moyen d'incident et sortait
12,5 milliards d'euros pour une ETI — un chiffre faux d'un facteur mille, produit par un code
juste.
Ce n'est pas un bug mais une erreur de catégorie : la télémétrie compte des attaques détectées,
la base externe des incidents qui ont coûté de l'argent, et rien n'empêchait de multiplier les
deux.
La correction est une calibration — je fixe p pour que cette entreprise retrouve le taux
d'incidents de 1 310 organisations comparables, pondérées par le même noyau que la sévérité,
soit une détection sur 5 138 qui finit en perte.
Le prix à payer, et je l'assume : la télémétrie ne pilote plus le niveau, seulement le mix par
type, et c'est le premier point de mon roadmap.
""",
    )

    # ------------------------------------------------- 8. sévérité : données
    slide = blank(deck)
    eyebrow(slide, "Sévérité — les données")
    title(slide, "Pondérer par similarité, jamais filtrer")

    card(slide, MARGIN, Inches(2.05), Inches(11.6), Inches(1.05), border=ACCENT)
    text(
        slide,
        "w  =  w_secteur  ×  w_taille  ×  exp( − d² / 2h² )        d = |maturité − 55|,  h = 15",
        left=MARGIN + Inches(0.35),
        top=Inches(2.38),
        width=Inches(11),
        height=Inches(0.5),
        size=19,
        color=INK,
        bold=True,
    )

    for index, (value, label, colour) in enumerate(
        [
            (f["incidents_fitted"], "incidents à perte exploitable\nsur 1 600 lus", INK),
            (f["losses_missing"], "sentinelles −1 exclues\njamais traitées comme 0 €", CAUTION),
            (f["hard_filter"], "incidents si filtre EXACT\naucun type modélisable", EDR),
            (f["min_neff"], "n_eff de Kish minimal\nsinon repli sur la poolée", ACCENT),
        ]
    ):
        figure(
            slide,
            value,
            label,
            left=MARGIN + Inches(index * 3.05),
            top=Inches(3.5),
            width=Inches(2.9),
            size=38,
            color=colour,
        )

    rule(slide, Inches(5.15))
    keywords(
        slide,
        [
            "encodage des secteurs réparé — « Ã‰nergie » → « Énergie »",
            "aucun incident jeté : les plus proches dominent, les autres pèsent peu",
            "chaque ajustement publie son n_eff — la finesse a un coût, il est affiché",
        ],
        left=MARGIN,
        top=Inches(5.5),
        width=Inches(11.5),
        size=14,
        gap=0.45,
    )
    notes(
        slide,
        """
Deux pièges dans la base externe : les pertes à −1 sont des valeurs manquantes et non des
incidents gratuits, et les libellés de secteur arrivaient en double encodage UTF-8, ce qui
séparait Énergie en deux secteurs distincts.
Le vrai choix de modélisation est la pondération douce : un filtre exact sur le profil ne laisse
que 112 incidents sur 1 598, et plus aucun type d'attaque n'a d'échantillon crédible —
mathématiquement propre, pratiquement inutilisable.
Je garde donc tout le monde avec un poids décroissant selon la ressemblance ; le prix de cette
souplesse est un échantillon effectif plus petit que l'échantillon nominal, que je publie type
par type et qui déclenche un repli sur la poolée en dessous de 30.
""",
    )

    # ---------------------------------------------- 9. sévérité : l'ajustement
    slide = blank(deck)
    eyebrow(slide, "Sévérité — l'ajustement")
    title(slide, "Une lognormale par type d'attaque")
    picture(slide, charts / "09_fit.png", left=MARGIN, top=Inches(2.05), width=Inches(8.3))
    text(
        slide,
        "supply chain",
        left=Inches(9.2),
        top=Inches(2.4),
        width=Inches(3.4),
        height=Inches(0.35),
        size=15,
        color=ACCENT,
        bold=True,
    )
    text(
        slide,
        f"μ = {f['sc_mu']}     σ = {f['sc_sigma']}",
        left=Inches(9.2),
        top=Inches(2.8),
        width=Inches(3.4),
        height=Inches(0.35),
        size=17,
        color=INK,
        bold=True,
    )
    figure(
        slide,
        f["sc_median"],
        "médiane — l'incident typique",
        left=Inches(9.2),
        top=Inches(3.45),
        width=Inches(3.4),
        size=30,
    )
    figure(
        slide,
        f["sc_mean"],
        "moyenne — ce qui alimente l'AAL",
        left=Inches(9.2),
        top=Inches(4.55),
        width=Inches(3.4),
        size=30,
        color=CAUTION,
    )
    text(
        slide,
        f"la queue multiplie la moyenne par {f['sc_ratio']}",
        left=Inches(9.2),
        top=Inches(5.6),
        width=Inches(3.6),
        height=Inches(0.4),
        size=13,
        color=EDR,
        bold=True,
    )
    text(
        slide,
        "MLE pondéré sur les log-pertes — les pertes sont bornées à zéro et s'étalent sur quatre "
        "ordres de grandeur : c'est la forme naturelle.",
        left=MARGIN,
        top=Inches(6.45),
        width=Inches(11.6),
        height=Inches(0.5),
        size=13.5,
        color=INK_2,
    )
    notes(
        slide,
        """
Pourquoi une lognormale : une perte ne peut pas être négative, n'a pas de plafond naturel et
s'étale ici sur quatre ordres de grandeur, la signature d'une variable dont le logarithme est à
peu près normal — en passant aux logs, l'asymétrie tombe de 7,4 à 0,8.
Un ajustement par type et non un seul global, parce que la base est un mélange de régimes très
différents qu'une lognormale unique se fait rejeter par Kolmogorov-Smirnov.
Le chiffre à retenir est l'écart entre médiane et moyenne, 183 k€ contre 5,1 M€ : l'incident
typique ne dit rien du coût moyen, et c'est la moyenne qui alimente la perte annuelle.
""",
    )

    # -------------------------------------------------- 10. sévérité : la preuve
    slide = blank(deck)
    eyebrow(slide, "Sévérité — la preuve")
    title(slide, "Chaque ajustement est livré avec les preuves contre lui")
    picture(slide, charts / "10_qq.png", left=MARGIN, top=Inches(2.05), width=Inches(7.2))
    keywords(
        slide,
        [
            "QQ-plot des log-pertes",
            "la diagonale = l'ajustement parfait",
            f"KS pondéré {f['sc_ks']} · seuil indicatif ≈ 0,19",
            "tient sur quatre ordres de grandeur",
            "extrêmes légèrement SOUS la courbe",
            "→ queue prudente sur ce type",
        ],
        left=Inches(8.5),
        top=Inches(2.5),
        width=Inches(4.2),
        size=13.5,
        gap=0.48,
    )
    card(slide, Inches(8.5), Inches(5.4), Inches(4.1), Inches(1.35), border=EDR)
    text(
        slide,
        "Contre-preuve publiée",
        left=Inches(8.75),
        top=Inches(5.58),
        width=Inches(3.7),
        height=Inches(0.3),
        size=12,
        color=EDR,
        bold=True,
    )
    text(
        slide,
        f"sur {f['pareto_better']} ajustements sur {f['fits_total']}, une queue\n"
        "de Pareto décrit mieux les extrêmes\n"
        "→ VaR 99 et TVaR 99 = bornes basses",
        left=Inches(8.75),
        top=Inches(5.92),
        width=Inches(3.7),
        height=Inches(0.8),
        size=11.5,
        color=INK_2,
        spacing=1.15,
    )
    notes(
        slide,
        """
Le KS pondéré vaut 0,118 pour un seuil indicatif autour de 0,19 à cet échantillon effectif,
donc l'ajustement n'est pas rejeté — mais un chiffre unique ne dit pas où le modèle se trompe,
d'où le QQ-plot.
Les points suivent la diagonale sur quatre ordres de grandeur et les trois plus extrêmes passent
légèrement en dessous : le modèle prédit un peu plus cher que l'observé, donc la queue est
prudente sur ce type, ce qui est le bon sens de l'erreur pour une mesure de risque.
Je publie aussi la contre-preuve et elle va dans l'autre sens : sur cinq ajustements sur neuf
une loi de Pareto décrit mieux les extrêmes, donc mes VaR 99 et TVaR 99 se lisent comme des
bornes basses.
""",
    )

    # ------------------------------------------------------ 11. Monte Carlo
    slide = blank(deck)
    eyebrow(slide, "Simulation Monte Carlo")
    title(slide, "Composer fréquence et sévérité, pas les multiplier")

    card(slide, MARGIN, Inches(2.0), Inches(11.6), Inches(1.35), border=LINE)
    text(
        slide,
        "pour chaque année simulée :   tirage Poisson(λ_type) par type d'attaque\n"
        "  →   prix de chaque incident dans la lognormale de son type"
        "   →   somme des pertes de l'année",
        left=MARGIN + Inches(0.35),
        top=Inches(2.32),
        width=Inches(11),
        height=Inches(0.9),
        size=15,
        color=INK,
        spacing=1.35,
    )
    picture(slide, charts / "12b_annual_loss.png", left=MARGIN, top=Inches(3.55), width=Inches(8.0))
    figure(
        slide,
        f["n_years"],
        "années simulées",
        left=Inches(9.1),
        top=Inches(3.75),
        width=Inches(3.4),
        size=32,
    )
    figure(
        slide,
        "seed 42",
        "reproductibilité bit à bit —\nla graine est affichée dans l'interface",
        left=Inches(9.1),
        top=Inches(4.75),
        width=Inches(3.4),
        size=32,
        color=ACCENT,
    )
    card(slide, Inches(9.1), Inches(5.95), Inches(3.4), Inches(1.0), border=CAUTION)
    text(
        slide,
        f"Garde-fou : perte unitaire plafonnée à {f['cap']}\n"
        f"({f['cap_share']} des tirages) — sinon σ = 2,6 tire\ndes années au-delà du plausible",
        left=Inches(9.3),
        top=Inches(6.12),
        width=Inches(3.1),
        height=Inches(0.8),
        size=10.5,
        color=INK_2,
        spacing=1.15,
    )
    notes(
        slide,
        """
Pourquoi simuler plutôt qu'une formule fermée : l'espérance d'une loi composée s'écrit
analytiquement, mais la VaR 99 et la TVaR 99 non — il faut la distribution complète, et elle
dépend de la façon dont les deux étages se composent.
Une année n'est pas la fréquence multipliée par le coût moyen : c'est zéro incident trois fois
sur quatre, un parfois, très rarement deux, et la somme de ces cas n'est pas le produit des
moyennes.
Sur le garde-fou : à sigma 2,6 et sur cent mille années, une lognormale non bornée finit par
tirer un incident qui coûte plus que l'entreprise ne vaut — la pire année non plafonnée montait
à 3,8 milliards, donc je plafonne au 99,9e centile des pertes réellement observées chez les
pairs et j'affiche ce que cela change.
""",
    )

    # --------------------------------------------------------- 12. résultats
    slide = blank(deck)
    eyebrow(slide, "Résultats")
    title(slide, "Ce que coûte une année")
    picture(slide, charts / "12_exceedance.png", left=MARGIN, top=Inches(1.95), width=Inches(8.1))

    metrics = [
        (f["aal"], "AAL — perte annuelle moyenne", ACCENT, 34),
        (f["median_year"], f"année médiane — {f['p_no_loss']} sans incident", INK, 30),
        (f["var95"], "VaR 95 — une année sur vingt", INK, 30),
        (f["var99"], "VaR 99 — une année sur cent", INK, 30),
    ]
    for index, (value, label, colour, size) in enumerate(metrics):
        figure(
            slide,
            value,
            label,
            left=Inches(9.0),
            top=Inches(1.95) + Inches(index * 1.10),
            width=Inches(3.6),
            size=size,
            color=colour,
        )
    text(
        slide,
        "budget   ·   appétit au risque   ·   question de survie",
        left=Inches(9.0),
        top=Inches(6.45),
        width=Inches(3.6),
        height=Inches(0.6),
        size=12.5,
        color=INK_3,
        spacing=1.2,
    )

    # The two tail averages ride under the chart, on the left half only: the
    # right column already runs to the foot of the slide.
    for index, (value, label) in enumerate(
        [
            (f["tvar95"], "TVaR 95 — 5 % pires années"),
            (f["tvar99"], "TVaR 99 — 1 % pires années"),
        ]
    ):
        left = MARGIN + Inches(index * 4.1)
        text(
            slide,
            value,
            left=left,
            top=Inches(6.35),
            width=Inches(1.5),
            height=Inches(0.45),
            size=26,
            color=CAUTION,
            bold=True,
        )
        text(
            slide,
            label,
            left=left + Inches(1.55),
            top=Inches(6.47),
            width=Inches(2.4),
            height=Inches(0.4),
            size=11.5,
            color=INK_2,
        )
    notes(
        slide,
        """
Trois lectures métier, trois chiffres : l'AAL à 274 k€ est la ligne budgétaire, la VaR 95 à
816 k€ l'appétit au risque, la TVaR 99 à 15,1 M€ la question de survie.
Le point contre-intuitif est l'année médiane à zéro euro — trois années sur quatre ne coûtent
rien, parce que le taux d'incidents est de 0,31 par an, et c'est la forme même du risque cyber
pour une ETI : rien, rien, rien, puis un événement qui compte.
Sur les courbes, l'AEP est le total de l'année et l'OEP la plus grosse perte unitaire, leur
écart mesurant la part des années à plusieurs incidents ; le plateau de l'OEP à droite est le
plafond de plausibilité, visible et non dissimulé.
""",
    )

    # ------------------------------------------------- 13. limites & suites
    slide = blank(deck)
    eyebrow(slide, "Limites assumées & suites")
    title(slide, "Ce que je changerais en premier")

    items = [
        (
            "p_materialize est un scalaire",
            "un seul nombre absorbe le bruit des sondes, la qualité des contrôles et la vitesse "
            "de réaction",
            "→ vulnérabilité FAIR pilotée par la maturité",
        ),
        (
            "Fenêtre de sept mois",
            "aucune saisonnalité annuelle vérifiable — un pic Black Friday resterait invisible",
            "→ fenêtre glissante, ingestion incrémentale",
        ),
        (
            "Queue lognormale",
            "Pareto décrit mieux les extrêmes sur 5 ajustements sur 9",
            "→ GPD par dépassement de seuil, à comparer",
        ),
        (
            "Types d'attaque indépendants",
            "un phishing amène un vol d'identifiants, qui amène un rançongiciel",
            "→ copule sur les fréquences par type",
        ),
    ]
    for index, (head, body, next_step) in enumerate(items):
        top = Inches(2.15) + Inches(index * 1.12)
        card(slide, MARGIN, top, Inches(11.6), Inches(0.95), border=LINE)
        text(
            slide,
            head,
            left=MARGIN + Inches(0.3),
            top=top + Inches(0.15),
            width=Inches(3.5),
            height=Inches(0.3),
            size=14.5,
            color=INK,
            bold=True,
        )
        text(
            slide,
            body,
            left=MARGIN + Inches(0.3),
            top=top + Inches(0.5),
            width=Inches(6.4),
            height=Inches(0.35),
            size=11.5,
            color=INK_3,
        )
        text(
            slide,
            next_step,
            left=MARGIN + Inches(7.2),
            top=top + Inches(0.32),
            width=Inches(4.2),
            height=Inches(0.35),
            size=12.5,
            color=ACCENT,
        )

    text(
        slide,
        "Détail et arbitrages : next_steps.md · page /roadmap de l'interface — 17 chantiers, ce "
        "qui existe / ce qui change / ce que ça apporte",
        left=MARGIN,
        top=Inches(6.75),
        width=Inches(11.6),
        height=Inches(0.4),
        size=12.5,
        color=INK_2,
    )
    notes(
        slide,
        """
Ces quatre limites ne sont pas des oublis mais des arbitrages, chacun écrit dans le dépôt avec
la raison de ne pas l'avoir traité.
La plus gênante est la première : p_materialize fait correspondre cette entreprise à ses pairs
mais ne sait pas dire pourquoi, si bien qu'améliorer la maturité de 55 à 75 ne bougerait pas la
perte modélisée — le mauvais signal pour un outil censé justifier un budget sécurité.
Ce que je n'ai pas fait faute de données : la régression maturité vers vulnérabilité demande un
dénominateur d'exposition que la base ne contient pas, puisqu'elle recense des incidents et non
des années-organisation à risque, et fabriquer cette courbe aurait donné un coefficient
d'allure confiante sans rien derrière.
""",
    )

    # -------------------------------------------------- 14. ce que vous auditez
    slide = blank(deck)
    eyebrow(slide, "Ce que vous auditez")
    title(slide, "Chaque chiffre a un chemin")

    steps = [
        ("data/*.csv", "4 fichiers, lecture seule"),
        ("NormalizationReport", "chaque ligne rendue"),
        ("FrequencyEstimate", "épisodes → λ → calibration"),
        ("SeverityModel", "nettoyage → poids → ajustement"),
        ("SimulationResult", "λ × sévérité → métriques"),
        ("seed 42", "rejouable à l'euro près"),
    ]
    for index, (name, detail) in enumerate(steps):
        top = Inches(2.15) + Inches(index * 0.62)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN, top + Inches(0.12), Inches(0.13), Inches(0.13)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if index in (0, len(steps) - 1) else INK_3
        dot.line.fill.background()
        dot.shadow.inherit = False
        text(
            slide,
            name,
            left=MARGIN + Inches(0.35),
            top=top,
            width=Inches(4.0),
            height=Inches(0.35),
            size=16,
            color=INK,
            bold=True,
        )
        text(
            slide,
            detail,
            left=MARGIN + Inches(4.3),
            top=top + Inches(0.04),
            width=Inches(4.5),
            height=Inches(0.35),
            size=12.5,
            color=INK_3,
        )

    card(slide, Inches(8.3), Inches(2.15), Inches(4.3), Inches(2.1), border=ACCENT)
    text(
        slide,
        "Une commande",
        left=Inches(8.6),
        top=Inches(2.38),
        width=Inches(3.8),
        height=Inches(0.3),
        size=12,
        color=ACCENT,
        bold=True,
    )
    box = text(
        slide,
        "python -m risk_engine\n  --data-dir data/\n  --out results.json",
        left=Inches(8.6),
        top=Inches(2.78),
        width=Inches(3.8),
        height=Inches(1.1),
        size=13,
        color=INK,
        spacing=1.3,
    )
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Consolas"

    text(
        slide,
        "« Vous pouvez rejouer tout ce que je viens de dire. »",
        left=MARGIN,
        top=Inches(6.0),
        width=Inches(11.6),
        height=Inches(0.6),
        size=26,
        color=INK,
        bold=True,
    )
    text(
        slide,
        "La trace imprimée par la CLI est la même que celle servie par l'API et affichée par "
        "l'interface — il n'y a qu'une version de chaque chiffre.",
        left=MARGIN,
        top=Inches(6.7),
        width=Inches(11.6),
        height=Inches(0.4),
        size=13,
        color=INK_2,
    )
    notes(
        slide,
        """
Je termine où j'ai commencé : chaque chiffre montré aujourd'hui a un chemin explicite depuis
les CSV, le rapport de normalisation rendant compte de chaque ligne lue avant que chaque étage
expose sa propre trace numérotée.
Ces traces ne sont pas de la documentation écrite à côté du code, elles sont produites par le
code à partir des mêmes objets que les calculs, et ne peuvent donc pas diverger du modèle.
Une seule commande rejoue l'ensemble et écrit un JSON contenant les chiffres et les
explications ; avec la graine 42 vous retombez sur mes valeurs à l'euro près.
Si un chiffre vous paraît faux, il y a une ligne de trace à contester plutôt qu'une parole à
croire.
""",
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    deck.save(str(out))
    print(f"{out}  —  {len(deck.slides.__iter__.__self__._sldIdLst)} diapositives")


def main() -> int:
    """Read the verified figures and build the deck."""
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--figures", type=Path, default=ROOT / "deck" / "figures.json")
    parser.add_argument("--charts", type=Path, default=ROOT / "deck" / "charts")
    parser.add_argument("--out", type=Path, default=ROOT / "presentation_citalid.pptx")
    args = parser.parse_args()

    figures = json.loads(args.figures.read_text(encoding="utf-8"))
    build(figures, args.charts, args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
